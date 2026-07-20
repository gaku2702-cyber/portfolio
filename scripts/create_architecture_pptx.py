"""ポートフォリオサイトの仕組みを説明する PowerPoint を生成する。"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
OUT = ROOT / "ポートフォリオサイト仕組み説明.pptx"

# 16:9
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

NAVY = RGBColor(0x1B, 0x32, 0x64)
AMBER = RGBColor(0xD9, 0x77, 0x06)
DARK = RGBColor(0x23, 0x23, 0x23)
MUTED = RGBColor(0x52, 0x52, 0x52)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF7, 0xF8, 0xFA)
BOX_BG = RGBColor(0xEE, 0xF1, 0xF6)
LINE = RGBColor(0xC5, 0xCB, 0xD6)


def set_run(run, size=18, bold=False, color=DARK, font="Meiryo UI"):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_text_box(slide, left, top, width, height, text, size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run(run, size=size, bold=bold, color=color)
    return box


def add_bullets(slide, left, top, width, height, items, size=18, color=DARK, spacing=Pt(10)):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = spacing
        p.level = 0
        run = p.add_run()
        run.text = f"・{item}"
        set_run(run, size=size, color=color)
    return box


def add_rect(slide, left, top, width, height, fill=BOX_BG, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
    # softer corners
    try:
        shape.adjustments[0] = 0.08
    except Exception:
        pass
    return shape


def add_box_label(slide, left, top, width, height, text, fill=BOX_BG, text_color=DARK, size=14, bold=True):
    shape = add_rect(slide, left, top, width, height, fill=fill)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.anchor = MSO_ANCHOR.MIDDLE
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(0)
        p.space_after = Pt(0)
        run = p.add_run()
        run.text = line
        set_run(run, size=size, bold=bold, color=text_color)
    return shape


def add_arrow(slide, x1, y1, x2, y2, color=NAVY):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    conn.line.color.rgb = color
    conn.line.width = Pt(2)
    return conn


def add_footer(slide, page: int, total: int = 11):
    add_text_box(
        slide,
        Inches(0.6),
        Inches(7.1),
        Inches(10),
        Inches(0.3),
        "ポートフォリオサイトの仕組み",
        size=11,
        color=MUTED,
    )
    add_text_box(
        slide,
        Inches(11.5),
        Inches(7.1),
        Inches(1.2),
        Inches(0.3),
        f"{page} / {total}",
        size=11,
        color=MUTED,
        align=PP_ALIGN.RIGHT,
    )


def add_title_bar(slide, title, subtitle=None):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(1.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    add_text_box(slide, Inches(0.6), Inches(0.28), Inches(12), Inches(0.5), title, size=28, bold=True, color=WHITE)
    if subtitle:
        add_text_box(slide, Inches(0.6), Inches(0.72), Inches(12), Inches(0.35), subtitle, size=13, color=RGBColor(0xD0, 0xD7, 0xE5))
    # accent underline
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.15), SLIDE_W, Inches(0.06))
    accent.fill.solid()
    accent.fill.fore_color.rgb = AMBER
    accent.line.fill.background()


def blank_slide(prs):
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)


def slide_cover(prs):
    slide = blank_slide(prs)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(4.85), SLIDE_W, Inches(0.08))
    accent.fill.solid()
    accent.fill.fore_color.rgb = AMBER
    accent.line.fill.background()

    add_text_box(slide, Inches(0.9), Inches(2.2), Inches(11.5), Inches(0.5), "ARCHITECTURE OVERVIEW", size=14, color=AMBER, bold=True)
    add_text_box(slide, Inches(0.9), Inches(2.8), Inches(11.5), Inches(1), "ポートフォリオサイトの仕組み", size=40, bold=True, color=WHITE)
    add_text_box(
        slide,
        Inches(0.9),
        Inches(4.1),
        Inches(11.5),
        Inches(0.5),
        "静的 HTML × GitHub Pages × Google Apps Script",
        size=18,
        color=RGBColor(0xD0, 0xD7, 0xE5),
    )
    add_text_box(
        slide,
        Inches(0.9),
        Inches(5.3),
        Inches(11.5),
        Inches(0.4),
        "https://gaku2702-cyber.github.io/portfolio/",
        size=16,
        color=WHITE,
    )
    add_text_box(slide, Inches(0.9), Inches(6.5), Inches(11.5), Inches(0.3), "1 / 11", size=12, color=RGBColor(0xA8, 0xB3, 0xC7))


def slide_summary(prs):
    slide = blank_slide(prs)
    add_title_bar(slide, "一言で言うと", "フレームワークなしの静的マルチページサイト")
    cards = [
        ("静的 HTML", "本体は index.html\nデモは兄弟 HTML\nビルド不要"),
        ("GitHub Pages", "main へ push で公開\nファイル＝そのまま URL\n数分で反映"),
        ("GAS 問い合わせ", "フォームだけ裏側あり\nシート記録 + メール\npostMessage で結果"),
    ]
    x = Inches(0.7)
    for title, body in cards:
        add_rect(slide, x, Inches(1.8), Inches(3.8), Inches(4.2), fill=BOX_BG)
        top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(1.8), Inches(3.8), Inches(0.7))
        top.fill.solid()
        top.fill.fore_color.rgb = NAVY
        top.line.fill.background()
        add_text_box(slide, x + Inches(0.25), Inches(1.95), Inches(3.3), Inches(0.45), title, size=20, bold=True, color=WHITE)
        add_text_box(slide, x + Inches(0.3), Inches(2.8), Inches(3.2), Inches(2.8), body, size=18, color=DARK)
        x += Inches(4.05)
    add_footer(slide, 2)


def slide_stack(prs):
    slide = blank_slide(prs)
    add_title_bar(slide, "技術スタック", "モダン JS フレームワークは使っていない")
    rows = [
        ("マークアップ", "素の HTML（lang=\"ja\"）"),
        ("スタイル", "Tailwind CDN + index.html 内の CSS"),
        ("動き", "各 HTML 内のバニラ JavaScript"),
        ("フォント", "Cormorant Garamond / Noto Sans JP / Inter"),
        ("公開", "GitHub Pages（main ブランチ）"),
        ("お問い合わせ", "Google Apps Script → スプレッドシート + Gmail"),
    ]
    y = Inches(1.55)
    for i, (label, value) in enumerate(rows):
        bg = WHITE if i % 2 == 0 else LIGHT_BG
        add_rect(slide, Inches(0.7), y, Inches(11.9), Inches(0.72), fill=bg, line=LINE)
        add_text_box(slide, Inches(0.95), y + Inches(0.18), Inches(3.2), Inches(0.4), label, size=16, bold=True, color=NAVY)
        add_text_box(slide, Inches(4.3), y + Inches(0.18), Inches(8), Inches(0.4), value, size=16, color=DARK)
        y += Inches(0.78)
    add_footer(slide, 3)


def slide_architecture(prs):
    slide = blank_slide(prs)
    add_title_bar(slide, "全体構成", "ブラウザから見たデータの流れ")

    # Row 1
    add_box_label(slide, Inches(0.6), Inches(1.7), Inches(2.4), Inches(0.9), "訪問者\nブラウザ", fill=NAVY, text_color=WHITE, size=14)
    add_box_label(slide, Inches(3.6), Inches(1.7), Inches(2.6), Inches(0.9), "GitHub Pages", fill=NAVY, text_color=WHITE, size=14)
    add_arrow(slide, Inches(3.05), Inches(2.15), Inches(3.55), Inches(2.15))

    # Row 2 - site parts
    add_box_label(slide, Inches(7.0), Inches(1.5), Inches(2.6), Inches(0.7), "index.html", fill=AMBER, text_color=WHITE, size=14)
    add_box_label(slide, Inches(10.0), Inches(1.5), Inches(2.7), Inches(0.7), "デモ HTML", fill=BOX_BG, text_color=DARK, size=14)
    add_box_label(slide, Inches(7.0), Inches(2.45), Inches(2.6), Inches(0.7), "assets/", fill=BOX_BG, text_color=DARK, size=14)
    add_box_label(slide, Inches(10.0), Inches(2.45), Inches(2.7), Inches(0.7), "covers / hero / docs", fill=BOX_BG, text_color=DARK, size=13)

    add_arrow(slide, Inches(6.25), Inches(2.15), Inches(6.95), Inches(2.15))

    # Contact flow
    add_text_box(slide, Inches(0.6), Inches(3.5), Inches(12), Inches(0.35), "お問い合わせの経路", size=16, bold=True, color=NAVY)

    add_box_label(slide, Inches(0.6), Inches(4.1), Inches(2.5), Inches(0.85), "Contact\nフォーム", fill=NAVY, text_color=WHITE, size=14)
    add_box_label(slide, Inches(3.7), Inches(4.1), Inches(2.6), Inches(0.85), "GAS\nウェブアプリ", fill=AMBER, text_color=WHITE, size=14)
    add_box_label(slide, Inches(6.9), Inches(4.1), Inches(2.6), Inches(0.85), "スプレッド\nシート記録", fill=BOX_BG, text_color=DARK, size=14)
    add_box_label(slide, Inches(10.1), Inches(4.1), Inches(2.5), Inches(0.85), "Gmail\n通知", fill=BOX_BG, text_color=DARK, size=14)

    add_arrow(slide, Inches(3.15), Inches(4.5), Inches(3.65), Inches(4.5))
    add_arrow(slide, Inches(6.35), Inches(4.5), Inches(6.85), Inches(4.5))
    add_arrow(slide, Inches(9.55), Inches(4.5), Inches(10.05), Inches(4.5))

    add_bullets(
        slide,
        Inches(0.6),
        Inches(5.3),
        Inches(12),
        Inches(1.4),
        [
            "Works カード → デモ HTML を別タブで開く / モーダル / PDF / Contact へ",
            "フォーム送信は hidden iframe + POST。成功/失敗は postMessage で受け取る",
        ],
        size=15,
    )
    add_footer(slide, 4)


def slide_pages(prs):
    slide = blank_slide(prs)
    add_title_bar(slide, "ページ構成", "ファイル名がそのまま URL になる")
    rows = [
        ("index.html", "本体。Hero → Concept → Works → Contact"),
        ("imposition-demo.html", "JFX200 面付シミュレーター"),
        ("photo-cut-demo.html", "写真切り出しデモ"),
        ("color-correct-demo.html", "カラー補正デモ"),
        ("pman-task-demo.html", "MIS タスク管理モック"),
        ("proposal-viewer.html", "PDF.js による提案書ビューア"),
    ]
    y = Inches(1.5)
    for i, (name, role) in enumerate(rows):
        bg = WHITE if i % 2 == 0 else LIGHT_BG
        add_rect(slide, Inches(0.7), y, Inches(11.9), Inches(0.72), fill=bg, line=LINE)
        add_text_box(slide, Inches(0.95), y + Inches(0.18), Inches(4.2), Inches(0.4), name, size=15, bold=True, color=NAVY)
        add_text_box(slide, Inches(5.3), y + Inches(0.18), Inches(7), Inches(0.4), role, size=15, color=DARK)
        y += Inches(0.78)
    add_text_box(
        slide,
        Inches(0.7),
        Inches(6.3),
        Inches(12),
        Inches(0.4),
        "ナビはハッシュアンカー（#concept / #services / #contact）。ルーターは使わない。",
        size=14,
        color=MUTED,
    )
    add_footer(slide, 5)


def slide_top_flow(prs):
    slide = blank_slide(prs)
    add_title_bar(slide, "トップページの流れ", "index.html のセクション構成")

    steps = ["Hero", "Concept", "Works", "Contact"]
    x = Inches(0.7)
    for i, step in enumerate(steps):
        add_box_label(slide, x, Inches(1.8), Inches(2.5), Inches(1.0), step, fill=NAVY if i == 0 else BOX_BG, text_color=WHITE if i == 0 else DARK, size=18)
        if i < len(steps) - 1:
            add_arrow(slide, x + Inches(2.55), Inches(2.3), x + Inches(2.95), Inches(2.3))
        x += Inches(3.1)

    add_bullets(
        slide,
        Inches(0.7),
        Inches(3.3),
        Inches(12),
        Inches(3.2),
        [
            "固定ナビ（モバイルはハンバーガー → ドロワー）",
            "Hero：ブランド + CTA。背景は assets/hero/",
            "Works：.works-grid に作品カードを手書き配置",
            "モーダル #tool-detail：Word→InDesign の詳細説明",
            "スクロール演出：IntersectionObserver で .reveal を表示",
        ],
        size=17,
    )
    add_footer(slide, 6)


def slide_content(prs):
    slide = blank_slide(prs)
    add_title_bar(slide, "作品データの持ち方", "CMS や作品用 JSON はない")

    add_rect(slide, Inches(0.7), Inches(1.6), Inches(5.8), Inches(4.6), fill=BOX_BG)
    add_text_box(slide, Inches(1.0), Inches(1.9), Inches(5.2), Inches(0.5), "ソース・オブ・トゥルース", size=18, bold=True, color=NAVY)
    add_bullets(
        slide,
        Inches(1.0),
        Inches(2.5),
        Inches(5.2),
        Inches(3.3),
        [
            "作品一覧 = index.html の HTML",
            "カバー画像 = assets/covers/",
            "カテゴリ + 日本語タイトルを直書き",
            "アクションはデモ / モーダル / PDF / Contact",
        ],
        size=16,
    )

    add_rect(slide, Inches(6.8), Inches(1.6), Inches(5.8), Inches(4.6), fill=LIGHT_BG)
    add_text_box(slide, Inches(7.1), Inches(1.9), Inches(5.2), Inches(0.5), "例外・補助", size=18, bold=True, color=NAVY)
    add_bullets(
        slide,
        Inches(7.1),
        Inches(2.5),
        Inches(5.2),
        Inches(3.3),
        [
            "pman-task-demo は mock JSON を読む",
            "scripts/generate_work_covers.py でカバー再生成可",
            "公開サイトが読むのは HTML と画像パスだけ",
        ],
        size=16,
    )
    add_footer(slide, 7)


def slide_contact(prs):
    slide = blank_slide(prs)
    add_title_bar(slide, "お問い合わせの流れ", "スパム対策付きで GAS に送る")

    steps = [
        ("1", "入力チェック\nハニーポット / 3秒 / CAPTCHA"),
        ("2", "iframe POST\nGAS_URL へ送信"),
        ("3", "GAS 処理\nシート記録 + メール"),
        ("4", "結果表示\npostMessage"),
    ]
    x = Inches(0.5)
    for num, text in steps:
        add_box_label(slide, x, Inches(1.7), Inches(0.55), Inches(0.55), num, fill=AMBER, text_color=WHITE, size=16)
        add_box_label(slide, x + Inches(0.65), Inches(1.65), Inches(2.35), Inches(1.35), text, fill=BOX_BG, text_color=DARK, size=13)
        if num != "4":
            add_arrow(slide, x + Inches(3.1), Inches(2.3), x + Inches(3.35), Inches(2.3))
        x += Inches(3.2)

    add_bullets(
        slide,
        Inches(0.7),
        Inches(3.5),
        Inches(12),
        Inches(2.8),
        [
            "サイト側：index.html の GAS_URL と handleSubmit",
            "サーバ側：gas/contact-form.gs（サニタイズ → シート → MailApp）",
            "宛先：smartdtp.studio.works@gmail.com",
            "URL を変えたら index.html の GAS_URL を更新して push",
            "トラブル時は gas/README.md（シートに行が増えるか確認）",
        ],
        size=16,
    )
    add_footer(slide, 8)


def slide_deploy(prs):
    slide = blank_slide(prs)
    add_title_bar(slide, "公開・更新フロー", "サイトにビルドステップはない")

    steps = [
        ("1. 編集", "HTML / 画像を直す"),
        ("2. 確認", "ローカルで表示チェック"),
        ("3. push", "git push origin main"),
        ("4. 反映", "Pages が数分で更新"),
    ]
    x = Inches(0.6)
    for title, body in steps:
        add_rect(slide, x, Inches(1.8), Inches(2.9), Inches(2.4), fill=BOX_BG)
        head = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(1.8), Inches(2.9), Inches(0.7))
        head.fill.solid()
        head.fill.fore_color.rgb = NAVY
        head.line.fill.background()
        add_text_box(slide, x + Inches(0.15), Inches(1.95), Inches(2.6), Inches(0.45), title, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.2), Inches(2.8), Inches(2.5), Inches(1.1), body, size=16, color=DARK, align=PP_ALIGN.CENTER)
        x += Inches(3.15)

    add_bullets(
        slide,
        Inches(0.7),
        Inches(4.6),
        Inches(12),
        Inches(1.8),
        [
            "公開 URL：https://gaku2702-cyber.github.io/portfolio/",
            "color_correct/ のデスクトップ配布（build.ps1 等）は別系統",
            "サイト用の CI / package.json / Next.js はない",
        ],
        size=16,
    )
    add_footer(slide, 9)


def slide_repo(prs):
    slide = blank_slide(prs)
    add_title_bar(slide, "リポジトリの周辺物", "サイト本体以外のフォルダ")

    rows = [
        ("assets/", "カバー・ヒーロー・デモ用データ・提案 PDF"),
        ("gas/", "問い合わせ用 GAS ソースとセットアップ手順"),
        ("scripts/", "カバー生成・手順書 DOCX / 本 PPTX 生成"),
        ("color_correct/", "デスクトップ製品本体（サイトのデモとは別）"),
        ("photoshop_sab/", "写真切り出しデモ用の PSD / サンプル"),
        ("dtp_settings.json", "ローカル DTP 設定（公開パイプライン外）"),
    ]
    y = Inches(1.5)
    for i, (path, role) in enumerate(rows):
        bg = WHITE if i % 2 == 0 else LIGHT_BG
        add_rect(slide, Inches(0.7), y, Inches(11.9), Inches(0.72), fill=bg, line=LINE)
        add_text_box(slide, Inches(0.95), y + Inches(0.18), Inches(3.8), Inches(0.4), path, size=15, bold=True, color=NAVY)
        add_text_box(slide, Inches(4.9), y + Inches(0.18), Inches(7.4), Inches(0.4), role, size=15, color=DARK)
        y += Inches(0.78)
    add_footer(slide, 10)


def slide_summary_howto(prs):
    slide = blank_slide(prs)
    add_title_bar(slide, "まとめ：触るときの目安", "何を直すとき、どのファイルか")

    rows = [
        ("見た目・作品追加", "index.html の Works カード + assets/covers/"),
        ("デモの挙動", "各 *-demo.html"),
        ("お問い合わせ不具合", "GAS_URL / GAS 再デプロイ / シート行"),
        ("提案 PDF", "assets/docs/ とリンク先"),
        ("この資料の再生成", "python scripts/create_architecture_pptx.py"),
    ]
    y = Inches(1.55)
    for i, (what, where) in enumerate(rows):
        bg = WHITE if i % 2 == 0 else LIGHT_BG
        add_rect(slide, Inches(0.7), y, Inches(11.9), Inches(0.8), fill=bg, line=LINE)
        add_text_box(slide, Inches(0.95), y + Inches(0.2), Inches(3.8), Inches(0.45), what, size=16, bold=True, color=NAVY)
        add_text_box(slide, Inches(4.9), y + Inches(0.2), Inches(7.4), Inches(0.45), where, size=16, color=DARK)
        y += Inches(0.88)

    add_text_box(
        slide,
        Inches(0.7),
        Inches(6.2),
        Inches(12),
        Inches(0.4),
        "編集 → 確認 → push。それだけが公開フロー。",
        size=16,
        bold=True,
        color=AMBER,
    )
    add_footer(slide, 11)


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_cover(prs)
    slide_summary(prs)
    slide_stack(prs)
    slide_architecture(prs)
    slide_pages(prs)
    slide_top_flow(prs)
    slide_content(prs)
    slide_contact(prs)
    slide_deploy(prs)
    slide_repo(prs)
    slide_summary_howto(prs)

    prs.save(OUT)
    print(f"Wrote {OUT} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
